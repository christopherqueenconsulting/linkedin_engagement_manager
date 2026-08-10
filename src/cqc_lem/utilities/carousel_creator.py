"""Carousel rendering + a structured, deterministic image-selection strategy for
content slides.

LinkedIn carousel imagery — research-backed principles encoded below (from Buffer /
PostNitro / Hootsuite analyses of high-engagement document posts):

- CONTENT (middle) slides should carry a relevant image so they don't read as
  text-on-blank. The image must REINFORCE the slide's single idea — a concept,
  metaphor, or data-viz cue — never be random stock filler unrelated to the point.
- Keep a CONSISTENT visual style across slides (calm, professional, uncluttered) so
  the deck reads as one system, not a scrapbook.
- AVOID busy / text-heavy images that fight the slide's own overlaid text for
  attention. Prefer clean subjects with breathing room.
- The COVER and CTA slides are their own thing (bold title / clear ask); we leave
  those layouts as-is and only enrich the middle slides.

Sourcing is deterministic and cheap-first: Pexels stock (a free API call) is the
default; Replicate generation is opt-in, low-rate, and only used with the user's
active avatar when a person likeness actually adds value (e.g. personal stories).
Every whether/source decision is seeded by (post_id, slide_index) so regeneration
is stable and unit-testable, and every step degrades gracefully to a text-only
layout rather than crashing carousel generation.
"""
import os
import random
import tempfile
import urllib.request
from datetime import datetime
from typing import Callable, Optional, Union

from lxml.etree import tostring
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml
from pptx.slide import Slide
from pydantic import BaseModel, Field, HttpUrl, StrictStr, conlist
from pydantic_extra_types.color import Color


# Generic slide model for reusability
class CarouselSlide(BaseModel):
    """One slide's content — the shape every template's slide type inherits.

    Renderers read `image_path` (a LOCAL file) and nothing reads `image_url`, so a slide whose only
    image is a URL is treated as having no image at all.
    """

    title: Optional[StrictStr] = Field(None, description="Title or heading of the slide")
    content: Optional[StrictStr] = Field(None, description="Main content of the slide", max_length=500)
    image_url: Optional[HttpUrl] = Field(None, description="URL to an image")
    image_path: Optional[StrictStr] = Field(None, description="Path to an image")


# Specific Carousel Templates

class EducationalContentSlide(CarouselSlide):
    """Slide of an educational deck — `CarouselSlide` unchanged.

    Typed separately only so one template can gain fields without touching the others.
    """

    # Custom Educational Content slides could have additional elements if needed
    pass


class EducationalContentCarousel(BaseModel):
    """The AWARENESS-stage deck: a cover, 1-4 tip/step slides, then a call to action.

    Both `run_content_plan.create_carousel_content` and the `POST /api/generate-carousel` preview
    route pick the carousel class from the buyer stage, and `ai_helper.generate_carousel_content`
    names these fields verbatim in its prompt — so renaming one changes BOTH what the model is asked
    to return and which builder `create_ppt` dispatches to.
    """

    cover: EducationalContentSlide = Field(..., description="Cover Slide: A bold title that clearly states the topic.")
    contents: conlist(EducationalContentSlide, min_length=1, max_length=4) = Field(...,
                                                                                   description="Content Slides: Each slide covers one tip or step with a combination of short text and relevant visuals.")
    call_to_action: EducationalContentSlide = Field(...,
                                                    description="Final Slide: Summarize key points and add a call to action.")


class CaseStudySlide(CarouselSlide):
    """Slide of a case-study deck — `CarouselSlide` unchanged.

    Typed separately only so one template can gain fields without touching the others.
    """

    # Custom case study slides can have additional constraints or attributes
    pass


class CaseStudyCarousel(BaseModel):
    """The CONSIDERATION-stage deck: challenge → solution → results, an optional testimonial, a CTA.

    Same stage-map and prompt contract as `EducationalContentCarousel`.
    """

    cover: CaseStudySlide = Field(...,
                                  description="Cover Slide: Title of the case study with the client’s name or project outcome.")
    challenge: CaseStudySlide = Field(..., description="Slide 2: Brief description of the problem faced by the client.")
    solution: CaseStudySlide = Field(..., description="Slide 3: Explanation of the approach or solution.")
    results: CaseStudySlide = Field(..., description="Slide 4: Highlight measurable results with data or visuals.")
    testimonial: Optional[CaseStudySlide] = Field(None,
                                                  description="Slide 5 (Optional): Include a client quote or feedback.")
    call_to_action: CaseStudySlide = Field(...,
                                           description="Final Slide: Encourage viewers to reach out for more information.")


class PersonalStorySlide(CarouselSlide):
    """Slide of a personal-story deck — `CarouselSlide` unchanged.

    Typed separately only so one template can gain fields without touching the others.
    """

    # Custom personal story slide if needed
    pass


class PersonalStoryCarousel(BaseModel):
    """A personal-story deck: cover, 1-3 moments of the journey, a takeaway, then a CTA.

    Only the `POST /api/generate-carousel` preview route selects it (stages `personal`/`story`); the
    30-day plan's stage map never does. It is also the one deck whose slides may carry the user's
    avatar likeness — `CAROUSEL_AVATAR_RELEVANT_TYPES` holds exactly its content type.
    """

    cover: PersonalStorySlide = Field(..., description="Cover Slide: Compelling title introducing the personal story.")
    story_slides: conlist(PersonalStorySlide, min_length=1, max_length=3) = Field(...,
                                                                                  description="Slides 2-4: Key moments in the journey.")
    takeaway: PersonalStorySlide = Field(..., description="Slide 5: Summary or lessons learned from the experience.")
    call_to_action: PersonalStorySlide = Field(...,
                                               description="Final Slide: Encourage viewers to share their own stories or connect with you for further discussion.")


class IndustryInsightSlide(CarouselSlide):
    """Slide of an industry-insight deck — `CarouselSlide` unchanged.

    Typed separately only so one template can gain fields without touching the others.
    """

    # Custom industry insight slide for trends and insights
    pass


class IndustryInsightsCarousel(BaseModel):
    """The FALLBACK deck: cover, 1-4 trends/insights, then a CTA.

    Every stage that is not awareness, consideration or decision lands here, in both stage maps — so
    an unrecognised stage still produces a valid deck rather than failing.
    """

    cover: IndustryInsightSlide = Field(...,
                                        description="Cover Slide: Title with an attention-grabbing phrase for industry insights.")
    insights: conlist(IndustryInsightSlide, min_length=1, max_length=4) = Field(...,
                                                                                description="Slides 2-5: Individual trends or insights with visuals.")
    call_to_action: IndustryInsightSlide = Field(...,
                                                 description="Final Slide: Summary and call-to-action for opinions on the trends.")


class EventRecapSlide(CarouselSlide):
    """Slide of an event-recap deck — `CarouselSlide` unchanged.

    Typed separately only so one template can gain fields without touching the others.
    """

    # Custom event recap slide if additional details are needed
    pass


class EventRecapCarousel(BaseModel):
    """An event-recap deck: title/date cover, 1-3 highlights, then a CTA.

    No buyer stage maps to it, so nothing in the app generates one today; `create_ppt` still renders
    one if it is handed one.
    """

    cover: EventRecapSlide = Field(..., description="cover Slide: Event title and date.")
    key_moments: conlist(EventRecapSlide, min_length=1, max_length=3) = Field(...,
                                                                              description="Slides 2-4: Key takeaways or highlights from the event.")
    call_to_action: EventRecapSlide = Field(...,
                                            description="Final Slide: Thank attendees and provide a call-to-action for future events or download additional resources.")


class TestimonialSlide(CarouselSlide):
    """A testimonial, adding the client attribution to the base slide shape.

    `create_ppt_testimonial_carousel` renders `content` as the quote and `client_name` as the
    attribution; no layout consumes `client_logo_url`.
    """

    client_name: str = Field(..., description="Name of the client providing the testimonial")
    client_logo_url: Optional[HttpUrl] = Field(None, description="URL to the client’s logo or photo")


class TestimonialCarousel(BaseModel):
    """A client-testimonial deck: cover, 1-3 quotes, then a CTA.

    Like `EventRecapCarousel` no stage maps to it, so it is renderable but never generated. It is
    also the ONE type `create_ppt` dispatches without `post_id`/`user_id`, so its slides get no
    deterministic image selection.
    """

    cover: CarouselSlide = Field(...,
                                 description="Cover Slide: Cover slide with a title like 'What Our Clients Are Saying'.")
    testimonials: conlist(TestimonialSlide, min_length=1, max_length=3) = Field(...,
                                                                                description="Slides 2-4: Individual testimonials with a quote, client name, and photo.")
    call_to_action: CarouselSlide = Field(...,
                                          description="Final Slide: Call-to-action to encourage viewers to reach out for similar results.")


class ProductDemoSlide(CarouselSlide):
    """Slide of a product-demo deck — `CarouselSlide` unchanged.

    Typed separately only so one template can gain fields without touching the others.
    """

    # Custom product demo slide if additional details are needed
    pass


class ProductDemoCarousel(BaseModel):
    """The DECISION-stage deck: cover, the headline feature, 1-2 further features, then a CTA."""

    cover: ProductDemoSlide = Field(...,
                                    description="Cover Slide: Introduction to the product with a compelling headline.")
    main_feature: ProductDemoSlide = Field(...,
                                           description="Slide 2: Highlight the main feature of the product with an image.")
    additional_features: conlist(ProductDemoSlide, min_length=1, max_length=2) = Field(...,
                                                                                       description="Slides 3-4: Additional features of the product.")
    call_to_action: ProductDemoSlide = Field(...,
                                             description="Final Slide: Call-to-action to learn more or sign up for a demo.")


class PowerPointThemeColors(BaseModel):
    """The theme colour slots `convert_ppt_theme_colors` writes into a saved .pptx.

    Every field name is used VERBATIM as the `<a:clrScheme>` child element in the XPath, so these
    must stay the OOXML names. A field left unset keeps whatever colour the design template shipped.
    """

    dk1: Optional[Color] = Field(None, description="RGB color code for the 1st dark color in the theme")
    lt1: Optional[Color] = Field(None, description="RGB color code for the 1st light color in the theme")
    dk2: Optional[Color] = Field(None, description="RGB color code for the 2nd dark color in the theme")
    lt2: Optional[Color] = Field(None, description="RGB color code for the 2nd light color in the theme")
    accent1: Optional[Color] = Field(None, description="RGB color code for the 1st accent color in the theme")
    accent2: Optional[Color] = Field(None, description="RGB color code for the 2nd accent color in the theme")
    accent3: Optional[Color] = Field(None, description="RGB color code for the 3rd accent color in the theme")
    accent4: Optional[Color] = Field(None, description="RGB color code for the 4th accent color in the theme")
    accent5: Optional[Color] = Field(None, description="RGB color code for the 5th accent color in the theme")
    accent6: Optional[Color] = Field(None, description="RGB color code for the 6th accent color in the theme")
    hlink: Optional[Color] = Field(None, description="RGB color code for the hyperlink color in the theme")
    folHlink: Optional[Color] = Field(None,
                                      description="RGB color code for the followed hyperlink color in the theme")


def create_ppt(ppt_name, carousel_data: Union[
    EducationalContentCarousel,
    CaseStudyCarousel,
    PersonalStoryCarousel,
    IndustryInsightsCarousel,
    EventRecapCarousel,
    TestimonialCarousel,
    ProductDemoCarousel],
               my_theme: PowerPointThemeColors = PowerPointThemeColors(**{"lt1": "e9d437", "dk2": "a89816"}),
               design_number: int = 1,
               post_id: Optional[int] = None,
               user_id: Optional[int] = None):
    """Render `carousel_data` into a .pptx under `generated_designs/` and return the saved path.

    The carousel's CONCRETE TYPE is the dispatch key. A model with no matching branch is saved as an
    untouched copy of the design template rather than raising — a new carousel type that forgets to
    add its branch here produces an empty deck, not an error.

    `design_number` selects `carousel_designs/Design-{n}.pptx`; theme colours are applied to the file
    AFTER it is saved. `post_id`/`user_id` are threaded through only to seed the deterministic
    per-slide image and layout selection, so re-rendering the same post reproduces the same deck.
    """
    current_dir = os.path.dirname(__file__)
    generated_dir = os.path.join(current_dir, "generated_designs")
    os.makedirs(generated_dir, exist_ok=True)
    design_path = os.path.join(current_dir, f"carousel_designs/Design-{design_number}.pptx")
    prs = Presentation(design_path)

    if isinstance(carousel_data, EducationalContentCarousel):
        # Handle EducationalContentCarousel
        prs = create_ppt_educational_content_carousel(prs, carousel_data, post_id=post_id, user_id=user_id)
        pass
    elif isinstance(carousel_data, CaseStudyCarousel):
        # Handle CaseStudyCarousel
        prs = create_ppt_case_study_carousel(prs, carousel_data, post_id=post_id, user_id=user_id)
        pass
    elif isinstance(carousel_data, PersonalStoryCarousel):
        prs = create_ppt_personal_story_carousel(prs, carousel_data, post_id=post_id, user_id=user_id)
    elif isinstance(carousel_data, IndustryInsightsCarousel):
        prs = create_ppt_industry_insights_carousel(prs, carousel_data, post_id=post_id, user_id=user_id)
    elif isinstance(carousel_data, EventRecapCarousel):
        prs = create_ppt_event_recap_carousel(prs, carousel_data, post_id=post_id, user_id=user_id)
    elif isinstance(carousel_data, TestimonialCarousel):
        prs = create_ppt_testimonial_carousel(prs, carousel_data)
    elif isinstance(carousel_data, ProductDemoCarousel):
        prs = create_ppt_product_demo_carousel(prs, carousel_data, post_id=post_id, user_id=user_id)

    file_path = os.path.join(generated_dir, f"{ppt_name}.pptx")
    prs.save(file_path)

    convert_ppt_theme_colors(file_path, my_theme)

    return f"{file_path}"


def get_default_image_path() -> str:
    """Path to the placeholder image shipped beside this module, used as the builders' `default_path`."""
    # Get the default image path local to this file
    file_dir = os.path.dirname(__file__)
    default_image_path = os.path.join(file_dir, "images/image.png")
    return default_image_path


def get_pexels_image_path(query: str, default_path: Optional[str] = None) -> Optional[str]:
    """Download a Pexels image matching *query* to a temp file and return its path.

    Falls back to *default_path* when PEXELS_API_KEY is absent or the request fails.
    Pass default_path=None to get None back on failure (callers that must NOT post a
    placeholder — e.g. the carousel poster — rely on this to flag the post instead).
    """
    try:
        from cqc_lem.utilities.pexels_helper import get_photo
        photo = get_photo(query)
        url = photo.medium  # medium-size JPEG is a good balance for slides
        # Pexels' image CDN 403s the default urllib User-Agent, so send a browser one.
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        with urllib.request.urlopen(req, timeout=20) as resp, open(tmp.name, "wb") as fh:
            fh.write(resp.read())
        return tmp.name
    except Exception:
        return default_path


# ── Structured, deterministic content-slide image selection ───────────────────
# Avatar (person likeness) only helps where the slide is about the author's own
# narrative — using it elsewhere is off-brand. Keep this narrow on purpose.
CAROUSEL_AVATAR_RELEVANT_TYPES: set[str] = {"personal_story"}

_QUERY_STOPWORDS: set[str] = {
    "the", "a", "an", "and", "or", "but", "to", "of", "for", "in", "on", "at",
    "with", "your", "you", "our", "we", "how", "why", "what", "when", "is", "are",
    "be", "this", "that", "it", "as", "from", "by", "will", "can", "do", "not",
    "more", "less", "than", "then", "so", "if", "up", "out", "into", "about",
    "tip", "tips", "step", "steps", "way", "ways",
}


def _seeded_unit(post_id: Optional[int], slide_index: int, purpose: str) -> float:
    """Deterministic float in [0, 1) keyed by (post_id, slide_index, purpose)."""
    return random.Random(f"{post_id}:{slide_index}:{purpose}").random()


def _heuristic_image_query(title: Optional[str], content: Optional[str],
                           content_type: Optional[str]) -> str:
    """Local (no-API) keyword extraction — concrete visual terms from slide text."""
    text = " ".join([t for t in (title, content) if t]).lower()
    words = [w.strip(".,:;!?\"'()[]") for w in text.split()]
    keywords = [w for w in words if len(w) > 3 and w not in _QUERY_STOPWORDS]
    if not keywords:
        return (content_type or "professional business").replace("_", " ")
    return " ".join(keywords[:3])


def derive_image_query(title: Optional[str], content: Optional[str],
                       content_type: Optional[str]) -> str:
    """Derive visual search keywords that capture the slide's MEANING (not its raw
    title). Uses an ``lem-simple`` LLM call when enabled; falls back to a local
    keyword heuristic on failure or when disabled.
    """
    from cqc_lem.utilities.env_constants import CAROUSEL_IMAGE_QUERY_LLM
    from cqc_lem.utilities.logger import log_debug, log_warning

    if not CAROUSEL_IMAGE_QUERY_LLM:
        return _heuristic_image_query(title, content, content_type)
    try:
        from cqc_lem.utilities.ai.client import client
        prompt = (
            "Give 2-4 concrete visual search keywords for a clean, professional stock "
            "photo that visually reinforces this LinkedIn carousel slide's single idea. "
            "Return ONLY the space-separated keywords, no punctuation, no quotes.\n\n"
            f"Title: {title or ''}\nBody: {content or ''}"
        )
        resp = client.chat.completions.create(
            model="lem-simple",
            messages=[{"role": "user", "content": prompt}],
        )
        raw = (resp.choices[0].message.content or "").strip()
        cleaned = " ".join(
            w.strip(".,:;!?\"'()[]") for w in raw.replace("\n", " ").split()
        ).strip()
        if cleaned:
            log_debug("Derived carousel image query", ai_model="lem-simple")
            return cleaned[:80]
        return _heuristic_image_query(title, content, content_type)
    except Exception as e:
        log_warning("Carousel image-query LLM failed, using heuristic", exc=e,
                    ai_model="lem-simple")
        return _heuristic_image_query(title, content, content_type)


def _should_include_slide_image(post_id: Optional[int], slide_index: int) -> bool:
    from cqc_lem.utilities.env_constants import CAROUSEL_IMAGE_RATE, CAROUSEL_IMAGES_ENABLED
    if not CAROUSEL_IMAGES_ENABLED:
        return False
    if CAROUSEL_IMAGE_RATE >= 1.0:
        return True
    if CAROUSEL_IMAGE_RATE <= 0.0:
        return False
    return _seeded_unit(post_id, slide_index, "include") < CAROUSEL_IMAGE_RATE


def _user_has_active_avatar(user_id: Optional[int], post_id: Optional[int] = None) -> bool:
    """Is the user's avatar usable for a CAROUSEL slide right now? (issue #744 guardrails)"""
    if user_id is None:
        return False
    from cqc_lem.utilities.avatar.guardrails import AVATAR_SURFACE_CAROUSEL, avatar_allowed_for
    return avatar_allowed_for(user_id, surface=AVATAR_SURFACE_CAROUSEL, post_id=post_id)


# A derived slide query is usually a THING, not a scene the author stands in. Only queries that
# read as a person get the avatar; everything else goes to base Flux / Pexels, because prepending
# the LoRA trigger word to "quarterly dashboard metrics" asked it to insert the user's face into a
# scene never written to contain a person (issue #744).
_PERSON_QUERY_TERMS: set[str] = {
    "person", "people", "team", "founder", "leader", "leadership", "colleague", "colleagues",
    "client", "clients", "customer", "customers", "candidate", "mentor", "manager", "employee",
    "employees", "speaker", "audience", "meeting", "interview", "conversation", "handshake",
    "portrait", "headshot", "founder's", "coworker", "coworkers", "presenter", "presentation",
    "teacher", "student", "engineer", "developer", "designer", "consultant", "recruiter",
}


def _query_depicts_person(query: Optional[str], content_type: Optional[str]) -> bool:
    """True when the slide's scene plausibly contains the author.

    ``personal_story`` slides are about the author by definition, so they qualify regardless of
    the derived keywords — that is precisely the narrative the avatar exists to illustrate.
    """
    if (content_type or "") in CAROUSEL_AVATAR_RELEVANT_TYPES:
        return True
    words = {w.strip(".,:;!?\"'()[]").lower() for w in (query or "").split()}
    return bool(words & _PERSON_QUERY_TERMS)


def _should_generate_with_replicate(post_id: Optional[int], slide_index: int,
                                    user_id: Optional[int], content_type: Optional[str]) -> bool:
    """Replicate generation is gated: globally enabled, avatar-relevant content type,
    a deterministic low-rate sample, AND the user's guardrails allow the avatar here.
    """
    from cqc_lem.utilities.env_constants import CAROUSEL_REPLICATE_ENABLED, CAROUSEL_REPLICATE_RATE
    if not CAROUSEL_REPLICATE_ENABLED or user_id is None:
        return False
    if (content_type or "") not in CAROUSEL_AVATAR_RELEVANT_TYPES:
        return False
    if _seeded_unit(post_id, slide_index, "source") >= CAROUSEL_REPLICATE_RATE:
        return False
    return _user_has_active_avatar(user_id, post_id)


def _generate_avatar_slide_image(query: str, user_id: int, post_id: Optional[int] = None,
                                 content_type: Optional[str] = None,
                                 title: Optional[str] = None,
                                 content: Optional[str] = None) -> Optional[str]:
    from cqc_lem.utilities.logger import log_warning
    try:
        from cqc_lem.utilities.ai.ai_helper import generate_post_image
        from cqc_lem.utilities.ai.image_brief import build_image_brief
        from cqc_lem.utilities.avatar.guardrails import AVATAR_SURFACE_CAROUSEL
        # Brief off the slide's ACTUAL text — the old keyword-bag prompt ("query, professional,
        # clean minimal background...") is exactly the generic filler this engine replaces.
        slide_text = "\n".join(t for t in (title, content) if t) or query
        brief = build_image_brief(slide_text, surface="carousel", ratio="1:1")
        path = generate_post_image(
            brief.prompt, user_id, surface=AVATAR_SURFACE_CAROUSEL,
            post_id=post_id,
            depicts_person=_query_depicts_person(query, content_type),
            focal_concept=brief.focal_concept)
        return path or None
    except Exception as e:
        log_warning("Carousel avatar image generation failed, falling back to Pexels",
                    exc=e, user_id=user_id, api_provider="replicate")
        return None


def select_slide_image(
    *,
    title: Optional[str],
    content: Optional[str],
    content_type: str,
    post_id: Optional[int],
    slide_index: int,
    user_id: Optional[int] = None,
    default_path: Optional[str] = None,
) -> Optional[str]:
    """Deterministically decide whether a content slide gets an image, derive its
    search query, and resolve a source (Pexels-first, then optional avatar-based
    Replicate generation). Returns a local image path or ``default_path``/None.

    Never raises — every failure degrades to ``default_path`` so carousel generation
    (and the caller's text-only layout fallback) is never blocked.
    """
    from cqc_lem.utilities.env_constants import CAROUSEL_PEXELS_ENABLED

    if not _should_include_slide_image(post_id, slide_index):
        return default_path

    query = derive_image_query(title, content, content_type)

    if _should_generate_with_replicate(post_id, slide_index, user_id, content_type):
        generated = _generate_avatar_slide_image(query, user_id, post_id, content_type,
                                                 title=title, content=content)
        if generated:
            return generated
        # fall through to Pexels on generation failure

    if CAROUSEL_PEXELS_ENABLED:
        pexels_path = get_pexels_image_path(query, default_path)
        if pexels_path:
            return pexels_path

    return default_path


def _content_layout_pools() -> tuple[list[Callable], list[Callable]]:
    """(image-capable, text-only) content layout pools. Image-capable layouts have a
    picture placeholder; text-only layouts render title+body with no image slot.
    """
    image_layouts = [create_one_column_text_layout_slide, create_one_column_text_1_layout_slide]
    text_layouts = [create_title_and_body_layout_slide, create_title_and_body_1_layout_slide]
    return image_layouts, text_layouts


def choose_content_layout(
    image_path: Optional[str],
    post_id: Optional[int],
    slide_index: int,
    image_layouts: Optional[list[Callable]] = None,
    text_layouts: Optional[list[Callable]] = None,
) -> Callable:
    """Route a content slide to a layout. When ``image_path`` is set the slide is
    guaranteed an IMAGE-CAPABLE layout (with a picture placeholder) so the fetched
    image never lands on a layout that silently discards it; otherwise a text-only
    layout is used. The choice is deterministic per (post_id, slide_index).
    """
    default_img, default_txt = _content_layout_pools()
    image_layouts = image_layouts or default_img
    text_layouts = text_layouts or default_txt
    rng = random.Random(f"{post_id}:{slide_index}:layout")
    return rng.choice(image_layouts if image_path else text_layouts)


def _insert_picture_into_placeholder(slide: Slide, placeholder, image_path: str) -> None:
    """Insert *image_path* into a placeholder.

    If *placeholder* is a ``PicturePlaceholder`` (has ``insert_picture``), use its
    native method. Otherwise fall back to ``slide.shapes.add_picture(...)`` positioned
    to the placeholder's frame so a non-picture placeholder never crashes the slide.
    """
    if hasattr(placeholder, "insert_picture"):
        placeholder.insert_picture(image_path)
        return

    slide.shapes.add_picture(
        image_path,
        placeholder.left,
        placeholder.top,
        placeholder.width,
        placeholder.height,
    )


def create_ppt_educational_content_carousel(prs: Presentation, carousel: EducationalContentCarousel,
                                            post_id: Optional[int] = None,
                                            user_id: Optional[int] = None) -> Presentation:
    """Create a PowerPoint presentation for Educational Content Carousel.

    Parameters:
    - prs: Presentation object to add slides to.
    - carousel: EducationalContentCarousel containing carouseldata.

    """
    # Get the default image path local to this file
    default_image_path = get_default_image_path()

    # Slide 1: Cover
    cover_layouts = [create_title_layout_slide, create_title_only_layout_slide,
                     create_title_and_body_layout_slide]
    cover_slide_args = {
        "prs": prs,
        "title": carousel.cover.title,
        "subtitle": carousel.cover.content,
        "body_text": carousel.cover.content
    }
    random.choice(cover_layouts)(**cover_slide_args)

    # Slide 2-5: Content/Tips — structured, deterministic image selection + layout
    # routing so image-bearing slides land on an image-capable layout.
    for slide_index, content in enumerate(carousel.contents, start=2):
        image_path = getattr(content, "image_path", None) or select_slide_image(
            title=content.title, content=content.content, content_type="educational",
            post_id=post_id, slide_index=slide_index, user_id=user_id,
            default_path=default_image_path,
        )
        layout_fn = choose_content_layout(image_path, post_id, slide_index)
        layout_fn(prs=prs, title=content.title, body_text=content.content, image_path=image_path)

    # Slide 6: Conclusion
    conclusion_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide,
                          create_caption_only_layout_slide]
    cta_image_path = getattr(carousel.call_to_action, "image_path", None) or get_pexels_image_path(
        carousel.call_to_action.title or "success", default_image_path
    )
    conclusion_slide_args = {
        "prs": prs,
        "title": carousel.call_to_action.title,
        "description": carousel.call_to_action.content,
        "subtitle": carousel.call_to_action.content,
        "image_path": cta_image_path,
    }
    random.choice(conclusion_layouts)(**conclusion_slide_args)

    return prs  # Return the presentation for further modifications or saving


def create_ppt_case_study_carousel(prs: Presentation, case_study_carousel: CaseStudyCarousel,
                                   post_id: Optional[int] = None,
                                   user_id: Optional[int] = None) -> Presentation:
    """Create a PowerPoint presentation for Case Study Carousel.

    Parameters:
    - prs: Presentation object to add slides to.
    - case_study_carousel: CaseStudyCarousel instance with structured content for each slide.
    """
    # Slide 1: Cover
    cover_layouts = [create_title_layout_slide, create_section_header_layout_slide, create_title_only_layout_slide]
    cover_kwargs = {
        'prs': prs,
        'percentage': getattr(case_study_carousel.cover, 'percentage', ''),
        'title': case_study_carousel.cover.title,
        'subtitle': getattr(case_study_carousel.cover, 'subtitle', case_study_carousel.cover.content)
    }
    random.choice(cover_layouts)(**cover_kwargs)

    default_image_path = get_default_image_path()

    # Slide 2: Challenge
    challenge_image = getattr(case_study_carousel.challenge, 'image_path', None) or select_slide_image(
        title=case_study_carousel.challenge.title, content=case_study_carousel.challenge.content,
        content_type="case_study", post_id=post_id, slide_index=2, user_id=user_id,
        default_path=default_image_path,
    )
    choose_content_layout(challenge_image, post_id, 2)(
        prs=prs, title=case_study_carousel.challenge.title,
        body_text=case_study_carousel.challenge.content, image_path=challenge_image,
    )

    # Slide 3: Solution
    solution_image = getattr(case_study_carousel.solution, 'image_path', None) or select_slide_image(
        title=case_study_carousel.solution.title, content=case_study_carousel.solution.content,
        content_type="case_study", post_id=post_id, slide_index=3, user_id=user_id,
        default_path=default_image_path,
    )
    choose_content_layout(solution_image, post_id, 3)(
        prs=prs, title=case_study_carousel.solution.title,
        body_text=case_study_carousel.solution.content, image_path=solution_image,
    )

    # Slide 4: Results — no image keeps the big-number data-viz layout available.
    results_image = getattr(case_study_carousel.results, 'image_path', None) or select_slide_image(
        title=case_study_carousel.results.title, content=case_study_carousel.results.content,
        content_type="case_study", post_id=post_id, slide_index=4, user_id=user_id,
        default_path=default_image_path,
    )
    choose_content_layout(
        results_image, post_id, 4,
        text_layouts=[create_big_number_layout_slide, create_title_and_body_1_layout_slide],
    )(
        prs=prs, title=case_study_carousel.results.title,
        body_text=case_study_carousel.results.content, image_path=results_image,
        big_number=getattr(case_study_carousel.results, 'big_number', ''),
        subtitle=getattr(case_study_carousel.results, 'subtitle', case_study_carousel.results.content),
    )

    # Slide 5: Testimonial
    testimonial_layouts = [create_caption_only_layout_slide, create_blank_1_1_layout_slide]
    testimonial_kwargs = {
        'prs': prs,
        'image_path': getattr(case_study_carousel.testimonial, 'image_path', get_default_image_path()),
        'title': getattr(case_study_carousel.testimonial, 'title', "Testimonial"),
        'quote': getattr(case_study_carousel.testimonial, 'quote', case_study_carousel.testimonial.content),
        'author': getattr(case_study_carousel.testimonial, 'author',
                          getattr(case_study_carousel.testimonial, 'title', "Happy Client"))
    }
    random.choice(testimonial_layouts)(**testimonial_kwargs)

    # Slide 6: Call to Action
    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    cta_kwargs = {
        'prs': prs,
        'title': case_study_carousel.call_to_action.title,
        'subtitle': case_study_carousel.call_to_action.content,
        'description': case_study_carousel.call_to_action.content
    }
    random.choice(cta_layouts)(**cta_kwargs)

    return prs  # Return the presentation for further modifications or saving


def create_ppt_personal_story_carousel(prs: Presentation, carousel: PersonalStoryCarousel,
                                       post_id: Optional[int] = None,
                                       user_id: Optional[int] = None) -> Presentation:
    """Append a personal-story deck's slides to `prs` and return it.

    Story slides are numbered from 2, and that index is half the seed for both `select_slide_image`
    and `choose_content_layout` — so a re-render of the same post reproduces the same images and
    layouts. The cover and CTA are picked with unseeded `random.choice` and do vary between runs.
    """
    default_image = get_default_image_path()

    cover_layouts = [create_title_layout_slide, create_section_header_layout_slide, create_title_only_layout_slide]
    random.choice(cover_layouts)(
        prs=prs, title=carousel.cover.title, subtitle=carousel.cover.content,
        percentage="", body_text=carousel.cover.content
    )

    for slide_index, slide_data in enumerate(carousel.story_slides, start=2):
        image_path = getattr(slide_data, "image_path", None) or select_slide_image(
            title=slide_data.title, content=slide_data.content, content_type="personal_story",
            post_id=post_id, slide_index=slide_index, user_id=user_id, default_path=default_image,
        )
        choose_content_layout(image_path, post_id, slide_index)(
            prs=prs, title=slide_data.title, body_text=slide_data.content, image_path=image_path,
        )

    create_title_and_body_1_layout_slide(
        prs=prs, title=carousel.takeaway.title, body_text=carousel.takeaway.content
    )

    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    random.choice(cta_layouts)(
        prs=prs, title=carousel.call_to_action.title,
        description=carousel.call_to_action.content, subtitle=carousel.call_to_action.content
    )
    return prs


def create_ppt_industry_insights_carousel(prs: Presentation, carousel: IndustryInsightsCarousel,
                                          post_id: Optional[int] = None,
                                          user_id: Optional[int] = None) -> Presentation:
    """Append an industry-insights deck's slides to `prs` and return it.

    Insight slides are numbered from 2, the index that (with `post_id`) seeds image and layout
    selection, so a re-render of the same post is stable. Cover and CTA use unseeded `random.choice`.
    """
    default_image = get_default_image_path()

    cover_layouts = [create_title_layout_slide, create_title_only_layout_slide]
    random.choice(cover_layouts)(
        prs=prs, title=carousel.cover.title, subtitle=carousel.cover.content
    )

    for slide_index, insight in enumerate(carousel.insights, start=2):
        image_path = getattr(insight, "image_path", None) or select_slide_image(
            title=insight.title, content=insight.content, content_type="industry_insights",
            post_id=post_id, slide_index=slide_index, user_id=user_id, default_path=default_image,
        )
        choose_content_layout(image_path, post_id, slide_index)(
            prs=prs, title=insight.title, body_text=insight.content, image_path=image_path,
        )

    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    random.choice(cta_layouts)(
        prs=prs, title=carousel.call_to_action.title,
        description=carousel.call_to_action.content, subtitle=carousel.call_to_action.content
    )
    return prs


def create_ppt_event_recap_carousel(prs: Presentation, carousel: EventRecapCarousel,
                                    post_id: Optional[int] = None,
                                    user_id: Optional[int] = None) -> Presentation:
    """Append an event-recap deck's slides to `prs` and return it.

    Highlight slides are numbered from 2, which with `post_id` seeds image and layout selection.
    Nothing in the app generates an `EventRecapCarousel` today, so this only runs if one is built by
    hand and handed to `create_ppt`.
    """
    default_image = get_default_image_path()

    cover_layouts = [create_title_layout_slide, create_section_header_layout_slide]
    random.choice(cover_layouts)(
        prs=prs, title=carousel.cover.title, subtitle=carousel.cover.content,
        percentage=""
    )

    for slide_index, moment in enumerate(carousel.key_moments, start=2):
        image_path = getattr(moment, "image_path", None) or select_slide_image(
            title=moment.title, content=moment.content, content_type="event_recap",
            post_id=post_id, slide_index=slide_index, user_id=user_id, default_path=default_image,
        )
        choose_content_layout(image_path, post_id, slide_index)(
            prs=prs, title=moment.title, body_text=moment.content, image_path=image_path,
        )

    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    random.choice(cta_layouts)(
        prs=prs, title=carousel.call_to_action.title,
        description=carousel.call_to_action.content, subtitle=carousel.call_to_action.content
    )
    return prs


def create_ppt_testimonial_carousel(prs: Presentation, carousel: TestimonialCarousel) -> Presentation:
    """Append a testimonial deck's slides to `prs` and return it.

    Alone among the builders it takes no `post_id`/`user_id`, so its slides get no image selection at
    all — every testimonial renders as quote + attribution on the blank layout.
    """
    cover_layouts = [create_title_layout_slide, create_title_only_layout_slide]
    random.choice(cover_layouts)(
        prs=prs, title=carousel.cover.title, subtitle=carousel.cover.content
    )

    for testimonial in carousel.testimonials:
        create_blank_1_1_layout_slide(
            prs=prs,
            quote=f'"{testimonial.content}"',
            author=f"— {testimonial.client_name}"
        )

    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    random.choice(cta_layouts)(
        prs=prs, title=carousel.call_to_action.title,
        description=carousel.call_to_action.content, subtitle=carousel.call_to_action.content
    )
    return prs


def create_ppt_product_demo_carousel(prs: Presentation, carousel: ProductDemoCarousel,
                                     post_id: Optional[int] = None,
                                     user_id: Optional[int] = None) -> Presentation:
    """Append a product-demo deck's slides to `prs` and return it.

    The headline feature is slide 2 and the additional features run from 3, so the (post_id,
    slide_index) seed holds a re-render on the same images and layouts.
    """
    default_image = get_default_image_path()

    cover_layouts = [create_title_layout_slide, create_title_only_layout_slide]
    random.choice(cover_layouts)(
        prs=prs, title=carousel.cover.title, subtitle=carousel.cover.content
    )

    main_image = getattr(carousel.main_feature, "image_path", None) or select_slide_image(
        title=carousel.main_feature.title, content=carousel.main_feature.content,
        content_type="product_demo", post_id=post_id, slide_index=2, user_id=user_id,
        default_path=default_image,
    )
    choose_content_layout(main_image, post_id, 2)(
        prs=prs, title=carousel.main_feature.title, body_text=carousel.main_feature.content,
        image_path=main_image,
    )

    for offset, feature in enumerate(carousel.additional_features):
        slide_index = 3 + offset
        image_path = getattr(feature, "image_path", None) or select_slide_image(
            title=feature.title, content=feature.content, content_type="product_demo",
            post_id=post_id, slide_index=slide_index, user_id=user_id, default_path=default_image,
        )
        choose_content_layout(image_path, post_id, slide_index)(
            prs=prs, title=feature.title, body_text=feature.content, image_path=image_path,
        )

    cta_layouts = [create_section_title_and_description_layout_slide, create_custom_3_1_layout_slide]
    random.choice(cta_layouts)(
        prs=prs, title=carousel.call_to_action.title,
        description=carousel.call_to_action.content, subtitle=carousel.call_to_action.content
    )
    return prs


def debug_slide(slide):
    """Print each shape on `slide` with its type and placeholder index (developer aid, no callers)."""
    for shape in slide.shapes:
        print(f"Shape: {shape.name}, Type: {shape.shape_type}, Placeholder: {shape.placeholder_format.idx}")


def convert_ppt_theme_colors(ppt_path, theme_colors: PowerPointThemeColors):
    """Rewrite a saved deck's theme colour scheme IN PLACE and re-save it.

    Only fields SET on `theme_colors` are touched; the rest keep the design template's own colours.
    python-pptx exposes no writable theme part, so the modified XML is assigned to its private
    `_blob`. Each field name goes straight into the `<a:clrScheme>` XPath, so a name with no matching
    element raises IndexError rather than being skipped.
    """
    # Load the presentation
    prs = Presentation(ppt_path)

    # Get the Slide Master
    slide_master = prs.slide_master
    slide_master_part = slide_master.part

    # Get the Theme and part
    theme_part = slide_master_part.part_related_by(RT.THEME)
    theme = parse_xml(theme_part.blob)  # theme here is an <a:theme> element

    # For each of the attributes in the PowerPointThemeColors model, find the corresponding XML element and update the color value
    for field_name, field_value in theme_colors:
        if field_value:
            color_element = theme.xpath(f'a:themeElements/a:clrScheme/a:{field_name}/a:srgbClr')[0]
            # print(f"{field_name} color before: {color_element.get('val')}")
            set_color = field_value.as_hex(format="long").replace("#", "")
            # print(f"{field_name} color after: {set_color}")
            color_element.set('val', set_color.encode('utf-8'))

    # Serialize the modified XML back to the theme part
    theme_part._blob = tostring(theme)

    # print(f"Blob After: {theme_part.blob}")

    # Save the presentation
    prs.save(ppt_path)


def set_ppt_theme_colors(ppt_path, theme_colors: dict = None):
    """Same in-place theme rewrite as `convert_ppt_theme_colors`, driven by a name→hex dict.

    Omitting `theme_colors` sets EVERY slot to white, which flattens the design's palette rather than
    leaving it alone — the opposite of the model-driven variant's partial update. No caller today.
    """
    # Load default theme colors
    if theme_colors is None:
        theme_colors = {
            "dk1": "ffffff",
            "lt1": "ffffff",
            "dk2": "ffffff",
            "lt2": "ffffff",
            "accent1": "ffffff",
            "accent2": "ffffff",
            "accent3": "ffffff",
            "accent4": "ffffff",
            "accent5": "ffffff",
            "accent6": "ffffff",
            "hlink": "ffffff",
            "folHlink": "ffffff"
        }

    # Load the presentation
    prs = Presentation(ppt_path)

    # Get the Slide Master
    slide_master = prs.slide_master
    slide_master_part = slide_master.part

    # Get the Theme and part
    theme_part = slide_master_part.part_related_by(RT.THEME)
    theme = parse_xml(theme_part.blob)  # theme here is an <a:theme> element

    # For each of the theme color names in the themes dict, find the corresponding XML element and update the color value
    for theme_color_name, theme_color_hex_value in theme_colors.items():
        if theme_color_name:
            color_element = theme.xpath(f'a:themeElements/a:clrScheme/a:{theme_color_name}/a:srgbClr')[0]
            # print(f"{theme_color_name} color before: {color_element.get('val')}")
            color_element.set('val', theme_color_hex_value.encode('utf-8'))
            # print(f"{theme_color_name} color after: {theme_color_hex_value}")

    # Serialize the modified XML back to the theme part
    theme_part._blob = tostring(theme)

    # Save the presentation
    prs.save(ppt_path)


def get_attr_gracefully(obj, attr):
    """Return `getattr(obj, attr)`, or None if the lookup raises. No caller in the tree today."""
    try:
        return getattr(obj, attr)
    except Exception as e:
        print(f"Error accessing attribute ({attr}): {e}")
    return None


def create_title_layout_slide(prs: Presentation, title: str, subtitle: str, **kwargs) -> Slide:
    """Create a PowerPoint slide with the TITLE layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text.
    - subtitle: Additional descriptive text.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[0]  # Assuming 0 index is the TITLE layout
    slide = prs.slides.add_slide(layout)

    # Populate title and subtitle placeholders
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]  # Assuming index 1 is for subtitle

    # Set the text content
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    return slide


def create_section_header_layout_slide(prs: Presentation, percentage: str, title: str, subtitle: str,
                                       **kwargs) -> Slide:
    """Create a PowerPoint slide with the SECTION_HEADER layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - percentage: The main percentage or metric to highlight.
    - title: The main title text.
    - subtitle: Additional descriptive text.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[1]  # Assuming 1 index is the SECTION_HEADER layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    percentage_placeholder = slide.placeholders[0]
    title_placeholder = slide.placeholders[1]
    subtitle_placeholder = slide.placeholders[2]

    # Set the text content
    percentage_placeholder.text = percentage
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    return slide


def debug_master_slide_placeholders_and_text(design_number: int = 1):
    """Print every layout, placeholder and placeholder text in a `Design-{n}.pptx` template.

    A developer aid for discovering which layout index and placeholder ids a design exposes, since
    the `create_*_layout_slide` builders address both by number.
    """
    current_dir = os.path.dirname(__file__)
    design_path = os.path.join(current_dir, f"carousel_designs/Design-{design_number}.pptx")
    prs = Presentation(design_path)
    slide_master = prs.slide_master
    for slide_layout in slide_master.slide_layouts:
        print(f"Slide Layout: {slide_layout.name}")
        for placeholder in slide_layout.placeholders:
            print(f"Placeholder: {placeholder.name}, Type: {placeholder.placeholder_format.idx}")
            for shape in slide_layout.shapes:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    print(f"Shape: {shape.name}, Type: {shape.shape_type}, Placeholder: {phf.idx}")
                    if phf.idx == placeholder.placeholder_format.idx:
                        print(f"Text: {shape.text}")
        print("----")


def create_title_and_body_layout_slide(prs: Presentation, title: str, body_text: str, **kwargs) -> Slide:
    """Create a PowerPoint slide with the TITLE_AND_BODY layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - body_text: The body text content, which could include bullet points or paragraphs.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[2]  # Assuming 2 index is the TITLE_AND_BODY layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    body_placeholder = slide.placeholders[1]

    # Set the text content
    title_placeholder.text = title
    body_placeholder.text = body_text

    return slide


def create_title_and_two_columns_layout_slide(prs: Presentation, title: str,
                                              left_column_title: str, left_column_subtitle: str,
                                              right_column_title: str, right_column_subtitle: str,
                                              **kwargs
                                              ) -> Slide:
    """Create a PowerPoint slide with the TITLE_AND_TWO_COLUMNS layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - left_column_title: Title for the left column.
    - left_column_subtitle: Subtitle text for the left column.
    - right_column_title: Title for the right column.
    - right_column_subtitle: Subtitle text for the right column.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[3]  # Assuming 3 index is the TITLE_AND_TWO_COLUMNS layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    left_column_title_placeholder = slide.placeholders[1]
    left_column_subtitle_placeholder = slide.placeholders[2]
    right_column_title_placeholder = slide.placeholders[3]
    right_column_subtitle_placeholder = slide.placeholders[4]

    # Set the text content
    title_placeholder.text = title
    left_column_title_placeholder.text = left_column_title
    left_column_subtitle_placeholder.text = left_column_subtitle
    right_column_title_placeholder.text = right_column_title
    right_column_subtitle_placeholder.text = right_column_subtitle

    return slide


def create_title_only_layout_slide(prs: Presentation, title: str, **kwargs) -> Slide:
    """Create a PowerPoint slide with the TITLE_ONLY layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[4]  # Assuming 4 index is the TITLE_ONLY layout
    slide = prs.slides.add_slide(layout)

    # Populate the title placeholder
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Return the slide object for further customization
    return slide


def create_main_point_layout_slide(prs: Presentation, title: str, **kwargs) -> Slide:
    """Create a PowerPoint slide with the MAIN_POINT layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide, typically a key point or highlight.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[6]  # Assuming 6 index is the MAIN_POINT layout
    slide = prs.slides.add_slide(layout)

    # Populate the title placeholder
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Return the slide object for further customization
    return slide


def create_one_column_text_layout_slide(prs: Presentation, title: str, body_text: str, image_path: str = None,
                                        **kwargs) -> Slide:
    """Create a PowerPoint slide with the ONE_COLUMN_TEXT layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - body_text: The body text content for the left column.
    - image_path (optional): Path to the image file to be inserted into the picture placeholder.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[5]  # Assuming 5 index is the ONE_COLUMN_TEXT layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    body_placeholder = slide.placeholders[1]
    picture_placeholder = slide.placeholders[2]

    # Set the text content
    title_placeholder.text = title
    body_placeholder.text = body_text

    # Insert the image if the path is provided
    if image_path:
        _insert_picture_into_placeholder(slide, picture_placeholder, image_path)

    # Return the slide object for further customization
    return slide


def create_section_title_and_description_layout_slide(prs: Presentation, title: str, description: str,
                                                      **kwargs) -> Slide:
    """Create a PowerPoint slide with the SECTION_TITLE_AND_DESCRIPTION layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - description: Subtitle or description text for additional context.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[7]  # Assuming 7 index is the SECTION_TITLE_AND_DESCRIPTION layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    description_placeholder = slide.placeholders[1]

    # Set the text content
    title_placeholder.text = title
    description_placeholder.text = description

    # Return the slide object for further customization
    return slide


def create_caption_only_layout_slide(prs: Presentation, title: str, image_path: str = None, **kwargs) -> Slide:
    """Create a PowerPoint slide with the CAPTION_ONLY layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The caption or title text for the slide.
    - image_path: Path to the image file to be inserted into the picture placeholder.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[8]  # Assuming 8 index is the CAPTION_ONLY layout
    slide = prs.slides.add_slide(layout)

    # Populate the title placeholder
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Insert the image if the path is provided
    picture_placeholder = slide.placeholders[2]
    if image_path:
        _insert_picture_into_placeholder(slide, picture_placeholder, image_path)

    # Return the slide object for further customization
    return slide


def create_big_number_layout_slide(prs: Presentation, big_number: str, subtitle: str, **kwargs) -> Slide:
    """Create a PowerPoint slide with the BIG_NUMBER layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - big_number: The main number or percentage to highlight on the slide.
    - subtitle: Subtitle text for additional context or explanation.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[9]  # Assuming 9 index is the BIG_NUMBER layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    big_number_placeholder = slide.placeholders[0]
    subtitle_placeholder = slide.placeholders[1]

    # Set the text content
    big_number_placeholder.text = big_number
    subtitle_placeholder.text = subtitle

    # Return the slide object for further customization
    return slide


def create_blank_layout_slide(prs: Presentation, **kwargs) -> Slide:
    """Create a PowerPoint slide with the BLANK layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing complete customization.
    """
    # Get the layout and add a blank slide
    layout = prs.slide_layouts[10]  # Assuming 10 index is the BLANK layout
    slide = prs.slides.add_slide(layout)

    # Return the slide object for full customization
    return slide


def create_custom_6_1_layout_slide(prs: Presentation, title: str, columns: list[dict], **kwargs) -> Slide:
    """Create a PowerPoint slide with the CUSTOM_6_1 (Title and Three Columns) layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - columns: A list of dictionaries, each containing 'metric', 'sub_title', and 'description' for each column.
               Example format: [{'metric': 'XX%', 'sub_title': 'Column 1 Title', 'description': 'Details for Column 1'}, ...]
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[11]  # Assuming 11 index is the CUSTOM_6_1 layout
    slide = prs.slides.add_slide(layout)

    # Set the main title
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Populate each column's placeholders
    for i, column_content in enumerate(columns):
        if i > 2:
            break  # Only three columns available

        metric_placeholder = slide.placeholders[3 + i * 3]
        sub_title_placeholder = slide.placeholders[4 + i * 3]
        description_placeholder = slide.placeholders[5 + i * 3]

        # Set content for each column
        metric_placeholder.text = column_content['metric']
        sub_title_placeholder.text = column_content['sub_title']
        description_placeholder.text = column_content['description']

    # Return the slide object for further customization
    return slide


def create_title_only_1_1_layout_slide(prs: Presentation, title: str, **kwargs) -> Slide:
    """Create a PowerPoint slide with the TITLE_ONLY_1_1 layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[12]  # Assuming 12 index is the TITLE_ONLY_1_1 layout
    slide = prs.slides.add_slide(layout)

    # Populate the title placeholder
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Return the slide object for further customization
    return slide


def create_one_column_text_1_layout_slide(prs: Presentation, title: str, body_text: str, image_path: str = None,
                                          **kwargs) -> Slide:
    """Create a PowerPoint slide with the ONE_COLUMN_TEXT_1 layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - body_text: The body text content for the right column.
    - image_path: Path to the image file to be inserted into the picture placeholder.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[13]  # Assuming 13 index is the ONE_COLUMN_TEXT_1 layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    body_placeholder = slide.placeholders[1]
    picture_placeholder = slide.placeholders[2]

    # Set the text content
    title_placeholder.text = title
    body_placeholder.text = body_text

    # Insert the image if the path is provided
    if image_path:
        _insert_picture_into_placeholder(slide, picture_placeholder, image_path)

    # Return the slide object for further customization
    return slide


def create_blank_1_1_layout_slide(prs: Presentation, quote: str, author: str, **kwargs) -> Slide:
    """Create a PowerPoint slide with the BLANK_1_1 (Quote) layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - quote: The main quote or message for the slide.
    - author: The author or source of the quote.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[14]  # Assuming 14 index is the BLANK_1_1 layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    quote_placeholder = slide.placeholders[0]
    author_placeholder = slide.placeholders[1]

    # Set the text content
    quote_placeholder.text = quote
    author_placeholder.text = author

    # Return the slide object for further customization
    return slide


def create_title_and_two_columns_1_layout_slide(prs: Presentation, title: str, column_1: dict,
                                                column_2: dict, **kwargs) -> Slide:
    """Create a PowerPoint slide with the TITLE_AND_TWO_COLUMNS_1 layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - column_1: A dictionary with 'sub_title' and 'text' for the left column.
               Example format: {'sub_title': 'Column 1 Title', 'text': 'Content for Column 1'}
    - column_2: A dictionary with 'sub_title' and 'text' for the right column.
               Example format: {'sub_title': 'Column 2 Title', 'text': 'Content for Column 2'}
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[15]  # Assuming 15 index is the TITLE_AND_TWO_COLUMNS_1 layout
    slide = prs.slides.add_slide(layout)

    # Set the main title
    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title

    # Populate left column's placeholders
    column_1_subtitle_placeholder = slide.placeholders[1]
    column_1_text_placeholder = slide.placeholders[2]
    column_1_subtitle_placeholder.text = column_1['sub_title']
    column_1_text_placeholder.text = column_1['text']

    # Populate right column's placeholders
    column_2_subtitle_placeholder = slide.placeholders[3]
    column_2_text_placeholder = slide.placeholders[4]
    column_2_subtitle_placeholder.text = column_2['sub_title']
    column_2_text_placeholder.text = column_2['text']

    # Return the slide object for further customization
    return slide


def create_title_and_body_1_layout_slide(prs: Presentation, title: str, body_text: str, **kwargs) -> Slide:
    """Create a PowerPoint slide with the TITLE_AND_BODY_1 layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide.
    - body_text: The body text content for the slide, supporting paragraphs or bullet points.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[16]  # Assuming 16 index is the TITLE_AND_BODY_1 layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    body_placeholder = slide.placeholders[1]

    # Set the text content
    title_placeholder.text = title
    body_placeholder.text = body_text

    # Return the slide object for further customization
    return slide


def create_custom_3_1_layout_slide(prs: Presentation, title: str, subtitle: str, **kwargs) -> Slide:
    """Create a PowerPoint slide with the CUSTOM_3_1 (Thanks) layout.

    Parameters:
    - prs: Presentation object to add slides to.
    - title: The main title text for the slide, typically a closing or thank-you message.
    - subtitle: Additional information or call-to-action text.
    - **kwargs: Additional keyword arguments are thrown away to allow for flexible function calls.

    Returns:
    - Slide: The PowerPoint slide object, allowing further customization.
    """
    # Get the layout and add a slide
    layout = prs.slide_layouts[17]  # Assuming 17 index is the CUSTOM_3_1 (Thanks) layout
    slide = prs.slides.add_slide(layout)

    # Populate the placeholders
    title_placeholder = slide.placeholders[0]
    subtitle_placeholder = slide.placeholders[1]

    # Set the text content
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    # Return the slide object for further customization
    return slide


def test_create_educational_ppt():
    """Test function to create an EducationalContentCarousel presentation and test the create_ppt function.
    """
    # Sample data for an educational content carousel
    carousel_data = {
        "cover": {
            "title": "5 Tips for Boosting Productivity",
            "content": "Practical ways to get more done every day"
        },
        "contents": [
            {"title": "Tip 1", "content": "Set clear goals for each day."},
            {"title": "Tip 2", "content": "Take regular breaks to recharge."},
            {"title": "Tip 3", "content": "Eliminate distractions to stay focused."},
            {"title": "Tip 4", "content": "Use productivity tools to track your progress."}
        ],
        "call_to_action": {
            "title": "Comment Below",
            "content": "Which Tip Will You Try First?"
        }
    }

    # PPT Name with tiemstamp suffix
    ppt_name = f"Educational_Carousel_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    my_ppt = create_ppt(ppt_name, EducationalContentCarousel(**carousel_data))
    print(f"Presentation created: {my_ppt}")


def test_create_case_study_ppt():
    """Test function to create an CaseStudyCarousel presentation and test the create_ppt function.
    """
    # Sample data for an case study content carousel
    carousel_data = {
        "cover": {
            "title": "Case Study: Successful Project",
            "content": "An in-depth look at our successful project with Client X"
        },
        "challenge": {
            "title": "The Challenge",
            "content": "Client X faced significant challenges in their market due to increased competition and changing customer preferences.",
            "image_path": get_default_image_path()
        },
        "solution": {
            "title": "Our Solution",
            "content": "We implemented a comprehensive strategy that included market analysis, customer engagement, and product innovation."
        },
        "results": {
            "title": "The Results",
            "content": "Our solution led to a 30% increase in market share and a 20% increase in customer satisfaction.",
            "image_path": get_default_image_path(),
            "big_number": "30%",
            "subtitle": "Increase in Market Share"
        },
        "testimonial": {
            "title": "Client Testimonial",
            "content": "Working with this team was a game-changer for our business. Their expertise and dedication were evident in every step of the process.",
            "image_path": get_default_image_path(),
            "quote": "This team transformed our business.",
            "author": "Jane Doe, CEO of Client X"
        },
        "call_to_action": {
            "title": "Get in Touch",
            "content": "Contact us to learn how we can help your business achieve similar results."
        }
    }

    # PPT Name with timestamp suffix
    ppt_name = f"Case_Study_Carousel_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    my_ppt = create_ppt(ppt_name, CaseStudyCarousel(**carousel_data))
    print(f"Presentation created: {my_ppt}")


def test_caption_only_slide(design_number: int = 1):
    """Render one CAPTION_ONLY slide into `generated_designs/` as a manual eyeball check.

    Despite the `test_` prefix this is not part of the pytest suite (which collects `tests/`) — it is
    run by hand from this module's `__main__` block and writes a real file.
    """
    current_dir = os.path.dirname(__file__)
    generated_dir = os.path.join(current_dir, "generated_designs")
    os.makedirs(generated_dir, exist_ok=True)
    design_path = os.path.join(current_dir, f"carousel_designs/Design-{design_number}.pptx")
    prs = Presentation(design_path)
    create_caption_only_layout_slide(prs, "Caption Only Slide", image_path=get_default_image_path())
    ppt_name = f"Caption Only Slide_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    file_path = os.path.join(generated_dir, f"{ppt_name}.pptx")
    prs.save(file_path)

    print(f"Created: {file_path}")


# ── Research-backed carousel templates ────────────────────────────────────────
# Each template is a dict of color and style parameters.
# Derived from Buffer/PostNitro/Hootsuite analysis of highest-engagement carousels.
CAROUSEL_TEMPLATES: dict[str, dict] = {
    "bold_listicle": {
        "label": "Bold Listicle",
        "description": "White slides, numbered badge circles, rainbow accents. Best for: tips, tools, mistakes.",
        "layout":      "listicle",
        "cover_bg":    (15, 23, 42),
        "cover_text":  (255, 255, 255),
        "cover_accent": (59, 130, 246),
        "content_bg":  (255, 255, 255),
        "title_color": (15, 23, 42),
        "body_color":  (71, 85, 105),
        "bottom_bar":  (15, 23, 42),
        "badge_colors": [(59, 130, 246), (16, 185, 129), (239, 68, 68), (139, 92, 246), (245, 158, 11), (14, 165, 233)],
    },
    "minimal_dark": {
        "label": "Minimal Dark",
        "description": "Black slides, huge left-aligned titles, gold accents. Best for: bold opinions, predictions.",
        "layout":      "dark_minimal",
        "cover_bg":    (10, 10, 10),
        "cover_text":  (255, 255, 255),
        "cover_accent": (251, 191, 36),
        "content_bg":  (18, 18, 18),
        "title_color": (255, 255, 255),
        "body_color":  (163, 163, 163),
        "bottom_bar":  (30, 30, 30),
        "badge_colors": [(251, 191, 36), (251, 146, 60), (52, 211, 153), (129, 140, 248), (248, 113, 113), (34, 211, 238)],
    },
    "stat_reveal": {
        "label": "Stat Reveal",
        "description": "Each slide title displayed ENORMOUS centered. Best for: data insights, research findings.",
        "layout":      "stat_big",
        "cover_bg":    (30, 64, 175),
        "cover_text":  (255, 255, 255),
        "cover_accent": (147, 197, 253),
        "content_bg":  (239, 246, 255),
        "title_color": (30, 64, 175),
        "body_color":  (55, 65, 81),
        "bottom_bar":  (30, 64, 175),
        "badge_colors": [(30, 64, 175)] * 6,
    },
    "step_framework": {
        "label": "Step Framework",
        "description": "Visual progress dots at top, arrow-bulleted body. Best for: how-to guides, playbooks.",
        "layout":      "step_progress",
        "cover_bg":    (4, 120, 87),
        "cover_text":  (255, 255, 255),
        "cover_accent": (110, 231, 183),
        "content_bg":  (255, 255, 255),
        "title_color": (6, 78, 59),
        "body_color":  (55, 65, 81),
        "bottom_bar":  (4, 120, 87),
        "badge_colors": [(16, 185, 129)] * 6,
    },
    "story_arc": {
        "label": "Story Arc",
        "description": "Cream slides, giant quote marks, square numbered badges. Best for: personal stories.",
        "layout":      "quote_pull",
        "cover_bg":    (120, 53, 15),
        "cover_text":  (255, 255, 255),
        "cover_accent": (253, 186, 116),
        "content_bg":  (255, 251, 235),
        "title_color": (92, 45, 0),
        "body_color":  (120, 53, 15),
        "bottom_bar":  (120, 53, 15),
        "badge_colors": [(245, 158, 11), (234, 88, 12), (217, 70, 239), (99, 102, 241), (239, 68, 68), (16, 185, 129)],
    },
}

DEFAULT_TEMPLATE = "bold_listicle"


def _wrap_text(text: str, font, max_px: int, draw) -> list[str]:
    """Greedy word-wrap `text` to lines no wider than `max_px` (measured via `draw`).

    A single token wider than `max_px` (long URL/word) is hard-broken into
    margin-fitting chunks so no line ever bleeds past the slide edge. `draw` is a
    Pillow ImageDraw whose ``textlength`` is used to measure — passed in so this is a
    pure function (unit-testable without the renderer closure).
    """
    if not text:
        return []

    def _fits(s: str) -> bool:
        return draw.textlength(s, font=font) <= max_px

    def _break_long(word: str) -> list[str]:
        chunks, cur = [], ""
        for ch in word:
            if cur and not _fits(cur + ch):
                chunks.append(cur)
                cur = ch
            else:
                cur += ch
        if cur:
            chunks.append(cur)
        return chunks

    words, lines, cur = text.split(), [], ""
    for word in words:
        for piece in ([word] if _fits(word) else _break_long(word)):
            test = (cur + " " + piece).strip()
            if _fits(test):
                cur = test
            else:
                if cur:
                    lines.append(cur)
                cur = piece
    if cur:
        lines.append(cur)
    return lines


def _fit_and_crop_image(image_path: str, target_w: int, target_h: int):
    """Cover-fit + center-crop an image to exactly (target_w, target_h) as RGB.

    Scales so the image fully covers the target box (no letterboxing), then crops
    the overflow from the center. Raises on any decode failure so callers can fall
    back to a text-only render.
    """
    from PIL import Image
    with Image.open(image_path) as src:
        rgb = src.convert("RGB")  # flattens transparency / palette onto RGB
    sw, sh = rgb.size
    if sw <= 0 or sh <= 0:
        raise ValueError("empty source image")
    scale = max(target_w / sw, target_h / sh)
    nw, nh = max(1, round(sw * scale)), max(1, round(sh * scale))
    resized = rgb.resize((nw, nh), Image.LANCZOS)
    left = (nw - target_w) // 2
    top = (nh - target_h) // 2
    return resized.crop((left, top, left + target_w, top + target_h))


def _carousel_content_type(carousel_data) -> str:
    """Map a carousel model instance to the content_type string used by the shared
    image-selection engine (query derivation + avatar relevance).
    """
    return {
        EducationalContentCarousel: "educational",
        CaseStudyCarousel: "case_study",
        PersonalStoryCarousel: "personal_story",
        IndustryInsightsCarousel: "industry_insights",
        EventRecapCarousel: "event_recap",
        TestimonialCarousel: "testimonial",
        ProductDemoCarousel: "product_demo",
    }.get(type(carousel_data), "professional")


def create_carousel_slide_images(
    carousel_data: Union[
        EducationalContentCarousel,
        CaseStudyCarousel,
        PersonalStoryCarousel,
        IndustryInsightsCarousel,
        EventRecapCarousel,
        TestimonialCarousel,
        ProductDemoCarousel,
    ],
    post_id: int,
    output_dir: Optional[str] = None,
    template: str = DEFAULT_TEMPLATE,
    user_id: Optional[int] = None,
    # Legacy params retained for backward compat but ignored
    bg_color: tuple = (26, 86, 219),
    accent_color: tuple = (255, 255, 255),
    secondary_bg: tuple = (15, 52, 142),
) -> list[str]:
    """Render carousel slides as 1080x1080 PNG images using Pillow.

    Creates one image per slide in output_dir (defaults to
    assets/images/carousel/{post_id}/). Returns a list of absolute image paths.

    ``template`` selects a visual style from CAROUSEL_TEMPLATES. Defaults to
    DEFAULT_TEMPLATE ("bold_listicle").

    CONTENT (middle) slides composite a relevant image into a bottom photo band via
    the shared, deterministic ``select_slide_image`` engine (Pexels-first, optional
    avatar-gated generation), seeded by (post_id, slide_index). Text reflows into the
    area above the band so nothing overlaps or clips. When no image is selected or the
    decode fails, the slide renders exactly as before (text-only). Cover + CTA slides
    are left as-is.
    """
    from PIL import Image, ImageDraw, ImageFont

    from cqc_lem.utilities.logger import log_warning

    W, H = 1080, 1080
    WHITE = (255, 255, 255)

    tmpl = CAROUSEL_TEMPLATES.get(template, CAROUSEL_TEMPLATES[DEFAULT_TEMPLATE])
    cover_bg     = tmpl["cover_bg"]
    cover_text   = tmpl["cover_text"]
    cover_accent = tmpl["cover_accent"]
    content_bg   = tmpl["content_bg"]
    title_color  = tmpl["title_color"]
    body_color   = tmpl["body_color"]
    bottom_bar   = tmpl["bottom_bar"]
    badge_colors = tmpl["badge_colors"]
    layout       = tmpl.get("layout", "listicle")

    if output_dir is None:
        current_dir = os.path.dirname(__file__)
        assets_root = os.path.join(current_dir, "..", "assets", "images", "carousel", str(post_id))
        output_dir = os.path.realpath(assets_root)
    os.makedirs(output_dir, exist_ok=True)

    # ── Font loader ───────────────────────────────────────────────────────────
    def _load_font(size: int, bold: bool = True) -> ImageFont.FreeTypeFont:
        bold_paths = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
            "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
            "/Library/Fonts/Arial Bold.ttf",
            "/System/Library/Fonts/Helvetica.ttc",
        ]
        reg_paths = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
            "/System/Library/Fonts/Supplemental/Arial.ttf",
            "/Library/Fonts/Arial.ttf",
            "/System/Library/Fonts/Helvetica.ttc",
        ]
        for path in (bold_paths if bold else reg_paths):
            if os.path.exists(path):
                try:
                    return ImageFont.truetype(path, size)
                except Exception:
                    continue
        return ImageFont.load_default(size=size)

    # ── Shared helpers ────────────────────────────────────────────────────────
    def _norm(text: str) -> str:
        for src, dst in {
            "—": "-", "–": "-", "‒": "-",
            "‘": "'", "’": "'",
            "“": '"', "”": '"',
            "…": "...", " ": " ",
            "•": "*", "→": "->", "←": "<-",
            "©": "(c)", "®": "(R)", "™": "(TM)",
        }.items():
            text = text.replace(src, dst)
        return text

    def _wrap(text: str, font, max_px: int, draw) -> list[str]:
        return _wrap_text(text, font, max_px, draw)

    def _block_h(lines, font, spacing, draw) -> int:
        total_h = 0
        for ln in lines:
            bb = draw.textbbox((0, 0), ln, font=font)
            total_h += (bb[3] - bb[1]) + spacing
        return total_h

    def _draw_block(draw, lines, font, x, y, fill, spacing=14,
                    centered=False, max_lines=99) -> int:
        for ln in lines[:max_lines]:
            bb = draw.textbbox((0, 0), ln, font=font)
            lw, lh = bb[2] - bb[0], bb[3] - bb[1]
            lx = (W - lw) // 2 if centered else x
            draw.text((lx, y), ln, font=font, fill=fill)
            y += lh + spacing
        return y

    def _rrect(draw, xy, radius, fill):
        draw.rounded_rectangle(list(xy), radius=radius, fill=fill)

    def _save(img: "Image.Image", idx: int) -> str:
        out = os.path.join(output_dir, f"slide_{idx:02d}.png")
        img.convert("RGB").save(out, "PNG", optimize=True)
        return out

    # ── Content-slide photo band ──────────────────────────────────────────────
    # A full-width band pinned above the slide footer. Preparing the panel up front
    # means a decode failure returns (None, None) and the slide renders text-only —
    # identical to today — instead of clipping text for an image that never lands.
    BAND_H = 360

    def _prep_band(image_path, footer_h, band_h=BAND_H):
        if not image_path:
            return None, None
        band_top = H - footer_h - band_h
        try:
            panel = _fit_and_crop_image(image_path, W, band_h)
        except Exception as e:
            log_warning("Carousel content image composite failed; rendering text-only",
                        exc=e, post_id=post_id)
            return None, None
        return panel, band_top

    def _place_band(img, draw, panel, band_top, accent):
        img.paste(panel, (0, band_top))
        draw.rectangle([(0, band_top), (W, band_top + 6)], fill=accent)

    # ══════════════════════════════════════════════════════════════════════════
    # LAYOUT: listicle  (Bold Listicle)
    # White content slides, colored numbered badge circle top-left, navy footer
    # ══════════════════════════════════════════════════════════════════════════
    def _listicle_cover(idx, total, title, body) -> str:
        f_t = _load_font(76, bold=True)
        f_s = _load_font(40, bold=False)
        f_l = _load_font(26, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        for row in range(H):
            draw.line([(0, row), (W, row)], fill=(*cover_accent, int(25 * (1 - row / H))))
        draw.rectangle([(0, 0), (10, H)], fill=cover_accent)
        pill = f"1 of {total}"
        pw = int(draw.textlength(pill, font=f_l)) + 36
        _rrect(draw, (50, 52, 50 + pw, 96), radius=22, fill=(*cover_accent, 200))
        draw.text((68, 60), pill, font=f_l, fill=cover_bg if sum(cover_accent) > 380 else WHITE)
        t_lines = _wrap(title, f_t, W - 140, draw)
        t_h = _block_h(t_lines[:4], f_t, 18, draw)
        t_y = max(150, (H // 2) - t_h // 2 - 80)
        y = _draw_block(draw, t_lines, f_t, 70, t_y, cover_text, spacing=18, centered=True, max_lines=4)
        draw.rectangle([(W // 2 - 50, y + 16), (W // 2 + 50, y + 22)], fill=cover_accent)
        s_lines = _wrap(body, f_s, W - 180, draw)
        _draw_block(draw, s_lines, f_s, 90, y + 46, fill=(*cover_text, 200), spacing=14, centered=True, max_lines=3)
        hint = "Swipe to read  >"
        hw = int(draw.textlength(hint, font=f_l))
        draw.text(((W - hw) // 2, H - 80), hint, font=f_l, fill=(*cover_accent, 160))
        return _save(img, idx)

    def _listicle_content(idx, total, title, body, badge_color, image_path=None) -> str:
        f_t = _load_font(64, bold=True)
        f_b = _load_font(38, bold=False)
        f_n = _load_font(50, bold=True)
        f_l = _load_font(26, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=content_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        BAR = 80
        panel, band_top = _prep_band(image_path, BAR)
        draw.rectangle([(0, 0), (W, 10)], fill=badge_color)
        cx, cy, cr = 118, 155, 68
        draw.ellipse([(cx - cr, cy - cr), (cx + cr, cy + cr)], fill=badge_color)
        num_str = str(idx - 1)
        nw = int(draw.textlength(num_str, font=f_n))
        bb = draw.textbbox((0, 0), num_str, font=f_n)
        nh = bb[3] - bb[1]
        draw.text((cx - nw // 2, cy - nh // 2 - 3), num_str, font=f_n, fill=WHITE)
        PAD = 62
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, cy + cr + 32, fill=title_color, spacing=12,
                        max_lines=2 if band_top else 3)
        draw.rectangle([(PAD, y + 14), (PAD + 90, y + 20)], fill=badge_color)
        y += 48
        b_lines = _wrap(body, f_b, W - PAD * 2, draw)
        _draw_block(draw, b_lines, f_b, PAD, y, fill=body_color, spacing=18,
                    max_lines=3 if band_top else 7)
        if panel is not None:
            _place_band(img, draw, panel, band_top, badge_color)
        draw.rectangle([(0, H - BAR), (W, H)], fill=bottom_bar)
        draw.rectangle([(0, H - BAR), (int(W * idx / total), H - BAR + 5)], fill=badge_color)
        cnt = f"{idx} / {total}"
        cw = int(draw.textlength(cnt, font=f_l))
        draw.text((W - cw - 36, H - 52), cnt, font=f_l, fill=(*badge_color, 210))
        return _save(img, idx)

    def _listicle_cta(idx, total, title, body) -> str:
        f_t = _load_font(72, bold=True)
        f_s = _load_font(40, bold=False)
        f_l = _load_font(26, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        for row in range(H):
            draw.line([(0, row), (W, row)], fill=(*cover_accent, int(25 * (row / H))))
        draw.rectangle([(W - 10, 0), (W, H)], fill=cover_accent)
        dx, dy, ds = 86, 96, 20
        draw.polygon([(dx, dy - ds), (dx + ds, dy), (dx, dy + ds), (dx - ds, dy)], fill=(*cover_accent, 180))
        pill = "Leave a comment below"
        pw = int(draw.textlength(pill, font=f_l)) + 36
        _rrect(draw, ((W - pw) // 2, 140, (W + pw) // 2, 184), radius=20, fill=(*cover_accent, 180))
        draw.text(((W - pw) // 2 + 18, 148), pill, font=f_l, fill=cover_bg if sum(cover_accent) > 380 else WHITE)
        cta_lines = _wrap(title, f_t, W - 140, draw)
        cta_h = _block_h(cta_lines[:3], f_t, 20, draw)
        cta_y = (H - cta_h) // 2 - 60
        y = _draw_block(draw, cta_lines, f_t, 70, cta_y, fill=cover_text, spacing=20, centered=True, max_lines=3)
        draw.rectangle([(W // 2 - 50, y + 18), (W // 2 + 50, y + 24)], fill=cover_accent)
        sub_lines = _wrap(body, f_s, W - 180, draw)
        _draw_block(draw, sub_lines, f_s, 90, y + 48, fill=(*cover_text, 190), spacing=16, centered=True, max_lines=3)
        draw.text((56, H - 76), f"{idx} / {total}", font=f_l, fill=(*cover_accent, 180))
        return _save(img, idx)

    # ══════════════════════════════════════════════════════════════════════════
    # LAYOUT: dark_minimal  (Minimal Dark)
    # Near-black slides; HUGE left-aligned title; thin gold rule; no badge circle
    # Left border bar is the only accent element on content slides
    # ══════════════════════════════════════════════════════════════════════════
    def _dark_cover(idx, total, title, body) -> str:
        f_t = _load_font(82, bold=True)
        f_s = _load_font(36, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Large decorative background character
        f_bg = _load_font(500, bold=True)
        draw.text((-40, H // 2 - 280), '"', font=f_bg, fill=(*cover_accent, 15))
        # Left thick accent bar
        draw.rectangle([(0, 0), (8, H)], fill=cover_accent)
        # Top-right slide counter
        cnt = f"01 / {total:02d}"
        cw = int(draw.textlength(cnt, font=f_l))
        draw.text((W - cw - 50, 52), cnt, font=f_l, fill=(*cover_accent, 160))
        # Title — large, left-aligned, starts at 30% from top
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD - 60, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, max(220, H // 3 - 60), cover_text, spacing=14, max_lines=4)
        # Thin gold rule
        draw.rectangle([(PAD, y + 24), (PAD + 120, y + 28)], fill=cover_accent)
        y += 52
        s_lines = _wrap(body, f_s, W - PAD - 60, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 170), spacing=16, max_lines=4)
        # Bottom "Swipe" hint
        hint = "SWIPE  >"
        draw.text((PAD, H - 80), hint, font=f_l, fill=(*cover_accent, 120))
        return _save(img, idx)

    def _dark_content(idx, total, title, body, badge_color, image_path=None) -> str:
        f_t = _load_font(70, bold=True)
        f_b = _load_font(36, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=content_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        panel, band_top = _prep_band(image_path, 60)
        # Left accent bar (thicker than cover)
        draw.rectangle([(0, 0), (16, H)], fill=badge_color)
        # Slide number — top right, muted
        cnt_str = f"{idx:02d} / {total:02d}"
        cw = int(draw.textlength(cnt_str, font=f_l))
        draw.text((W - cw - 50, 52), cnt_str, font=f_l, fill=(*badge_color, 140))
        # LARGE title left-aligned, starts high
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD - 60, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, 180, title_color, spacing=14,
                        max_lines=3 if band_top else 4)
        # Thin colored rule
        draw.rectangle([(PAD, y + 24), (PAD + 100, y + 27)], fill=badge_color)
        y += 52
        # Body text — smaller, muted
        b_lines = _wrap(body, f_b, W - PAD - 60, draw)
        _draw_block(draw, b_lines, f_b, PAD, y, fill=body_color, spacing=20,
                    max_lines=3 if band_top else 8)
        if panel is not None:
            _place_band(img, draw, panel, band_top, badge_color)
        # Bottom: thin line only
        draw.rectangle([(0, H - 60), (W, H - 58)], fill=(*badge_color, 80))
        return _save(img, idx)

    def _dark_cta(idx, total, title, body) -> str:
        f_t = _load_font(78, bold=True)
        f_s = _load_font(34, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        draw.rectangle([(0, 0), (8, H)], fill=cover_accent)
        cnt_str = f"{idx:02d} / {total:02d}"
        cw = int(draw.textlength(cnt_str, font=f_l))
        draw.text((W - cw - 50, 52), cnt_str, font=f_l, fill=(*cover_accent, 140))
        PAD = 70
        # CTA title — center of slide
        t_lines = _wrap(title, f_t, W - PAD - 60, draw)
        t_h = _block_h(t_lines[:3], f_t, 16, draw)
        t_y = (H - t_h) // 2 - 80
        y = _draw_block(draw, t_lines, f_t, PAD, t_y, cover_text, spacing=16, max_lines=3)
        draw.rectangle([(PAD, y + 24), (PAD + 100, y + 27)], fill=cover_accent)
        y += 52
        s_lines = _wrap(body, f_s, W - PAD - 60, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 170), spacing=18, max_lines=4)
        return _save(img, idx)

    # ══════════════════════════════════════════════════════════════════════════
    # LAYOUT: stat_big  (Stat Reveal)
    # Blue cover; content: title displayed ENORMOUS centered; tiny body below
    # ══════════════════════════════════════════════════════════════════════════
    def _stat_cover(idx, total, title, body) -> str:
        f_t = _load_font(74, bold=True)
        f_s = _load_font(38, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Wave-like two-tone split
        draw.rectangle([(0, H - 220), (W, H)], fill=(*cover_accent, 30))
        # Decorative large number "?" hinting at reveals
        f_deco = _load_font(300, bold=True)
        draw.text((W - 220, H // 2 - 200), "?", font=f_deco, fill=(*cover_accent, 18))
        # Slide counter pill top-center
        pill = f"1 of {total} reveals"
        pw = int(draw.textlength(pill, font=f_l)) + 36
        _rrect(draw, ((W - pw) // 2, 52, (W + pw) // 2, 96), radius=22, fill=(*cover_accent, 200))
        draw.text(((W - pw) // 2 + 18, 60), pill, font=f_l, fill=cover_bg)
        # Title centered
        PAD = 60
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        t_h = _block_h(t_lines[:4], f_t, 18, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, max(180, (H // 2) - t_h // 2 - 60), cover_text,
                        spacing=18, centered=True, max_lines=4)
        draw.rectangle([(W // 2 - 40, y + 20), (W // 2 + 40, y + 24)], fill=cover_accent)
        y += 48
        s_lines = _wrap(body, f_s, W - PAD * 2, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 200), spacing=14, centered=True, max_lines=3)
        return _save(img, idx)

    def _stat_content(idx, total, title, body, badge_color, image_path=None) -> str:
        f_huge = _load_font(100, bold=True)  # title as massive centered text
        f_body = _load_font(36, bold=False)
        f_l    = _load_font(24, bold=False)
        f_num  = _load_font(22, bold=True)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=content_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        panel, band_top = _prep_band(image_path, 60)
        # Top color band
        draw.rectangle([(0, 0), (W, 90)], fill=badge_color)
        # Step number in top band
        step_label = f"#{idx - 1}"
        draw.text((W // 2 - int(draw.textlength(step_label, font=f_num)) // 2, 30),
                  step_label, font=f_num, fill=(*content_bg, 220))
        PAD = 60
        # HUGE title — centered vertically in upper 65%; lifted to the top band when
        # an image occupies the lower third.
        title_max = 2 if band_top else 3
        t_lines = _wrap(title, f_huge, W - PAD * 2, draw)
        t_h = _block_h(t_lines[:title_max], f_huge, 18, draw)
        t_y = 140 if band_top else max(130, (H * 65 // 100) // 2 - t_h // 2)
        y = _draw_block(draw, t_lines, f_huge, PAD, t_y, title_color, spacing=18, centered=True,
                        max_lines=title_max)
        # Horizontal rule
        draw.rectangle([(PAD, y + 22), (W - PAD, y + 25)], fill=(*badge_color, 120))
        y += 50
        # Body small, centered
        b_lines = _wrap(body, f_body, W - PAD * 2, draw)
        _draw_block(draw, b_lines, f_body, PAD, y, fill=body_color, spacing=18, centered=True,
                    max_lines=2 if band_top else 5)
        if panel is not None:
            _place_band(img, draw, panel, band_top, badge_color)
        # Bottom strip
        draw.rectangle([(0, H - 60), (W, H)], fill=bottom_bar)
        progress_w = int(W * (idx - 1) / max(total - 2, 1))
        draw.rectangle([(0, H - 60), (progress_w, H - 55)], fill=badge_color)
        cnt = f"{idx - 1} / {total - 2}"
        cw = int(draw.textlength(cnt, font=f_l))
        draw.text((W - cw - 36, H - 42), cnt, font=f_l, fill=(*badge_color, 200))
        return _save(img, idx)

    def _stat_cta(idx, total, title, body) -> str:
        f_t = _load_font(74, bold=True)
        f_s = _load_font(38, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        draw.rectangle([(0, 0), (W, H)], fill=cover_bg)
        # Decorative corner triangle
        draw.polygon([(0, 0), (300, 0), (0, 300)], fill=(*cover_accent, 40))
        draw.polygon([(W, H), (W - 300, H), (W, H - 300)], fill=(*cover_accent, 40))
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        t_h = _block_h(t_lines[:3], f_t, 18, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, (H - t_h) // 2 - 80,
                        cover_text, spacing=18, centered=True, max_lines=3)
        draw.rectangle([(W // 2 - 40, y + 22), (W // 2 + 40, y + 25)], fill=cover_accent)
        y += 50
        s_lines = _wrap(body, f_s, W - PAD * 2, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 190), spacing=16, centered=True, max_lines=3)
        draw.text((W // 2 - int(draw.textlength(f"{idx} / {total}", font=f_l)) // 2, H - 70),
                  f"{idx} / {total}", font=f_l, fill=(*cover_accent, 160))
        return _save(img, idx)

    # ══════════════════════════════════════════════════════════════════════════
    # LAYOUT: step_progress  (Step Framework)
    # Green cover; content: top step-indicator strip, arrow-bulleted body
    # ══════════════════════════════════════════════════════════════════════════
    def _step_cover(idx, total, title, body) -> str:
        f_t = _load_font(76, bold=True)
        f_s = _load_font(38, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Diagonal accent block bottom-right
        draw.polygon([(W - 280, H), (W, H - 280), (W, H)], fill=(*cover_accent, 60))
        # Step label strip at top
        draw.rectangle([(0, 0), (W, 100)], fill=(*WHITE, 20))
        header = f"A {total - 2}-Step Framework"
        hw = int(draw.textlength(header, font=f_l))
        draw.text(((W - hw) // 2, 36), header, font=f_l, fill=(*cover_accent, 240))
        # Title centered
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        t_h = _block_h(t_lines[:4], f_t, 18, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, max(170, (H // 2) - t_h // 2 - 40),
                        cover_text, spacing=18, centered=True, max_lines=4)
        draw.rectangle([(PAD, y + 20), (W - PAD, y + 24)], fill=(*cover_accent, 180))
        y += 48
        s_lines = _wrap(body, f_s, W - PAD * 2, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 200), spacing=14, centered=True, max_lines=3)
        return _save(img, idx)

    def _step_content(idx, total, title, body, badge_color, image_path=None) -> str:
        f_t = _load_font(62, bold=True)
        f_b = _load_font(36, bold=False)
        f_n = _load_font(30, bold=True)
        title, body = _norm(title), _norm(body)
        content_steps = total - 2  # exclude cover + CTA
        step_num = max(1, idx - 1)

        img = Image.new("RGB", (W, H), color=content_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        panel, band_top = _prep_band(image_path, 60)

        # ── Step progress strip at top ────────────────────────────────────────
        STRIP = 100
        draw.rectangle([(0, 0), (W, STRIP)], fill=bottom_bar)
        dot_r = 18
        dot_gap = max(50, min(90, (W - 140) // max(content_steps, 1)))
        start_x = (W - (content_steps * dot_gap - (dot_gap - dot_r * 2))) // 2

        for i in range(content_steps):
            cx = start_x + i * dot_gap + dot_r
            cy = STRIP // 2
            if i + 1 < step_num:
                # completed — filled
                draw.ellipse([(cx - dot_r, cy - dot_r), (cx + dot_r, cy + dot_r)], fill=badge_color)
                check = "+"
                cw = int(draw.textlength(check, font=f_n))
                draw.text((cx - cw // 2, cy - 14), check, font=f_n, fill=WHITE)
            elif i + 1 == step_num:
                # current — filled with border
                draw.ellipse([(cx - dot_r - 4, cy - dot_r - 4), (cx + dot_r + 4, cy + dot_r + 4)],
                             fill=WHITE)
                draw.ellipse([(cx - dot_r, cy - dot_r), (cx + dot_r, cy + dot_r)], fill=badge_color)
                num = str(step_num)
                nw = int(draw.textlength(num, font=f_n))
                draw.text((cx - nw // 2, cy - 14), num, font=f_n, fill=WHITE)
            else:
                # future — outline only
                draw.ellipse([(cx - dot_r, cy - dot_r), (cx + dot_r, cy + dot_r)],
                             outline=(*badge_color, 100), width=3)
            # connecting line between dots
            if i < content_steps - 1:
                nx = start_x + (i + 1) * dot_gap + dot_r
                line_color = badge_color if i + 1 < step_num else (*badge_color, 60)
                draw.line([(cx + dot_r, cy), (nx - dot_r, cy)], fill=line_color, width=3)

        # ── Title ─────────────────────────────────────────────────────────────
        PAD = 62
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, STRIP + 40, title_color, spacing=12, max_lines=3)

        # Accent underline
        draw.rectangle([(PAD, y + 12), (PAD + 80, y + 17)], fill=badge_color)
        y += 42

        # ── Body with arrow bullets (fewer lines when a photo band is present) ──
        # Wrap to the width AFTER the arrow indent, keeping a right margin (PAD), so
        # bulleted lines never run past the slide edge.
        BULLET_INDENT = 52
        body_cap = 4 if band_top else 7
        for line_text in _wrap(body, f_b, W - (PAD + BULLET_INDENT) - PAD, draw)[:body_cap]:
            draw.text((PAD, y), "->", font=f_b, fill=badge_color)
            draw.text((PAD + BULLET_INDENT, y), line_text, font=f_b, fill=body_color)
            bb = draw.textbbox((0, 0), line_text, font=f_b)
            y += (bb[3] - bb[1]) + 20

        if panel is not None:
            _place_band(img, draw, panel, band_top, badge_color)

        # ── Bottom bar ────────────────────────────────────────────────────────
        BAR = 60
        draw.rectangle([(0, H - BAR), (W, H)], fill=bottom_bar)
        return _save(img, idx)

    def _step_cta(idx, total, title, body) -> str:
        f_t = _load_font(74, bold=True)
        f_s = _load_font(36, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        draw.polygon([(W - 280, H), (W, H - 280), (W, H)], fill=(*cover_accent, 60))
        draw.polygon([(0, 0), (280, 0), (0, 280)], fill=(*cover_accent, 40))
        # Checkmark large
        f_check = _load_font(120, bold=True)
        check_text = "Done!"
        cw = int(draw.textlength(check_text, font=f_check))
        draw.text(((W - cw) // 2, 120), check_text, font=f_check, fill=(*cover_accent, 220))
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD * 2, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, 380, cover_text, spacing=18, centered=True, max_lines=3)
        draw.rectangle([(PAD, y + 20), (W - PAD, y + 23)], fill=(*cover_accent, 180))
        y += 48
        s_lines = _wrap(body, f_s, W - PAD * 2, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 200), spacing=14, centered=True, max_lines=3)
        return _save(img, idx)

    # ══════════════════════════════════════════════════════════════════════════
    # LAYOUT: quote_pull  (Story Arc)
    # Cream content slides; giant quotation mark; editorial/magazine feel
    # Right amber border bar; title reads like a pull-quote
    # ══════════════════════════════════════════════════════════════════════════
    def _story_cover(idx, total, title, body) -> str:
        f_t = _load_font(74, bold=True)
        f_s = _load_font(38, bold=False)
        f_l = _load_font(24, bold=False)
        f_big = _load_font(260, bold=True)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Big decorative quote mark — watermark
        draw.text((50, 100), '"', font=f_big, fill=(*cover_accent, 35))
        # Warm texture: bottom band
        draw.rectangle([(0, H - 180), (W, H)], fill=(*cover_accent, 40))
        # Right border
        draw.rectangle([(W - 12, 0), (W, H)], fill=cover_accent)
        # Issue/series label top
        label = f"Part 1 of {total}"
        lw = int(draw.textlength(label, font=f_l))
        draw.text((W - lw - 36, 48), label, font=f_l, fill=(*cover_accent, 200))
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD - 80, draw)
        t_h = _block_h(t_lines[:4], f_t, 16, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, max(200, (H // 2) - t_h // 2 - 60),
                        cover_text, spacing=16, max_lines=4)
        # Amber rule
        draw.rectangle([(PAD, y + 20), (PAD + 100, y + 24)], fill=cover_accent)
        y += 50
        s_lines = _wrap(body, f_s, W - PAD - 80, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 190), spacing=16, max_lines=3)
        return _save(img, idx)

    def _story_content(idx, total, title, body, badge_color, image_path=None) -> str:
        f_quote = _load_font(180, bold=True)
        f_t     = _load_font(58, bold=True)
        f_b     = _load_font(36, bold=False)
        f_l     = _load_font(22, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=content_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        panel, band_top = _prep_band(image_path, 60)
        # Right accent border
        draw.rectangle([(W - 14, 0), (W, H)], fill=badge_color)
        # Slide number — top-left badge (square, not circle)
        draw.rectangle([(50, 50), (116, 116)], fill=badge_color)
        num_str = str(idx - 1)
        f_n = _load_font(42, bold=True)
        nw = int(draw.textlength(num_str, font=f_n))
        bb = draw.textbbox((0, 0), num_str, font=f_n)
        nh = bb[3] - bb[1]
        draw.text((83 - nw // 2, 83 - nh // 2 - 3), num_str, font=f_n, fill=WHITE)
        # Large decorative quote mark
        draw.text((50, 90), '"', font=f_quote, fill=(*badge_color, 25))
        # Title — larger, treated as pull-quote
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD - 80, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, 220, title_color, spacing=14,
                        max_lines=3 if band_top else 4)
        # Amber rule
        draw.rectangle([(PAD, y + 16), (PAD + 100, y + 20)], fill=badge_color)
        y += 46
        # Body text — softer
        b_lines = _wrap(body, f_b, W - PAD - 80, draw)
        _draw_block(draw, b_lines, f_b, PAD, y, fill=body_color, spacing=20,
                    max_lines=3 if band_top else 6)
        if panel is not None:
            _place_band(img, draw, panel, band_top, badge_color)
            # keep the signature right border visible over the photo band
            draw.rectangle([(W - 14, band_top), (W, H)], fill=badge_color)
        # Bottom: issue label
        draw.rectangle([(0, H - 60), (W - 14, H)], fill=bottom_bar)
        cnt = f"Part {idx - 1} of {total - 2}"
        draw.text((PAD, H - 44), cnt, font=f_l, fill=(*badge_color, 200))
        return _save(img, idx)

    def _story_cta(idx, total, title, body) -> str:
        f_quote = _load_font(180, bold=True)
        f_t = _load_font(68, bold=True)
        f_s = _load_font(36, bold=False)
        f_l = _load_font(24, bold=False)
        title, body = _norm(title), _norm(body)
        img = Image.new("RGB", (W, H), color=cover_bg)
        draw = ImageDraw.Draw(img, "RGBA")
        # Closing quote mark (right-aligned)
        draw.text((W - 200, H // 2 - 80), '"', font=f_quote, fill=(*cover_accent, 30))
        draw.rectangle([(W - 12, 0), (W, H)], fill=cover_accent)
        # "The End" style label
        label = "The Takeaway"
        lw = int(draw.textlength(label, font=f_l))
        draw.text((W - lw - 36, 48), label, font=f_l, fill=(*cover_accent, 200))
        PAD = 70
        t_lines = _wrap(title, f_t, W - PAD - 80, draw)
        t_h = _block_h(t_lines[:3], f_t, 16, draw)
        y = _draw_block(draw, t_lines, f_t, PAD, max(200, (H // 2) - t_h // 2 - 80),
                        cover_text, spacing=16, max_lines=3)
        draw.rectangle([(PAD, y + 20), (PAD + 100, y + 24)], fill=cover_accent)
        y += 50
        s_lines = _wrap(body, f_s, W - PAD - 80, draw)
        _draw_block(draw, s_lines, f_s, PAD, y, fill=(*cover_text, 190), spacing=16, max_lines=3)
        draw.text((PAD, H - 70), f"Part {idx} of {total}", font=f_l, fill=(*cover_accent, 160))
        return _save(img, idx)

    # ══════════════════════════════════════════════════════════════════════════
    # Dispatch to the right render functions
    # ══════════════════════════════════════════════════════════════════════════
    COVER_FN = {
        "listicle":      _listicle_cover,
        "dark_minimal":  _dark_cover,
        "stat_big":      _stat_cover,
        "step_progress": _step_cover,
        "quote_pull":    _story_cover,
    }
    CONTENT_FN = {
        "listicle":      _listicle_content,
        "dark_minimal":  _dark_content,
        "stat_big":      _stat_content,
        "step_progress": _step_content,
        "quote_pull":    _story_content,
    }
    CTA_FN = {
        "listicle":      _listicle_cta,
        "dark_minimal":  _dark_cta,
        "stat_big":      _stat_cta,
        "step_progress": _step_cta,
        "quote_pull":    _story_cta,
    }
    render_cover   = COVER_FN.get(layout, _listicle_cover)
    render_content = CONTENT_FN.get(layout, _listicle_content)
    render_cta     = CTA_FN.get(layout, _listicle_cta)

    # ── Collect slides from carousel model ────────────────────────────────────
    slides_data: list[tuple[str, str]] = []

    def _add(t, c):
        slides_data.append((t or "", c or ""))

    if isinstance(carousel_data, EducationalContentCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        for s in carousel_data.contents:
            _add(s.title, s.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, CaseStudyCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        _add(carousel_data.challenge.title, carousel_data.challenge.content)
        _add(carousel_data.solution.title, carousel_data.solution.content)
        _add(carousel_data.results.title, carousel_data.results.content)
        if carousel_data.testimonial:
            _add(carousel_data.testimonial.title, carousel_data.testimonial.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, PersonalStoryCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        for s in carousel_data.story_slides:
            _add(s.title, s.content)
        _add(carousel_data.takeaway.title, carousel_data.takeaway.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, IndustryInsightsCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        for s in carousel_data.insights:
            _add(s.title, s.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, EventRecapCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        for s in carousel_data.key_moments:
            _add(s.title, s.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, TestimonialCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        for s in carousel_data.testimonials:
            _add(s.title, s.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)
    elif isinstance(carousel_data, ProductDemoCarousel):
        _add(carousel_data.cover.title, carousel_data.cover.content)
        _add(carousel_data.main_feature.title, carousel_data.main_feature.content)
        for s in carousel_data.additional_features:
            _add(s.title, s.content)
        _add(carousel_data.call_to_action.title, carousel_data.call_to_action.content)

    # ── Render ────────────────────────────────────────────────────────────────
    content_type = _carousel_content_type(carousel_data)
    total = len(slides_data)
    image_paths = []
    for idx, (title, body) in enumerate(slides_data, start=1):
        if idx == 1:
            path = render_cover(idx, total, title, body)
        elif idx == total:
            path = render_cta(idx, total, title, body)
        else:
            bc = badge_colors[(idx - 2) % len(badge_colors)]
            # Shared deterministic engine; default_path=None so a miss => text-only
            # (never a placeholder image).
            slide_image = select_slide_image(
                title=title, content=body, content_type=content_type,
                post_id=post_id, slide_index=idx, user_id=user_id, default_path=None,
            )
            path = render_content(idx, total, title, body, bc, slide_image)
        image_paths.append(path)

    return image_paths


def create_carousel_pdf(image_paths: list[str], post_id: int,
                        output_dir: Optional[str] = None) -> Optional[str]:
    """Bundle already-rendered slide images into ONE multi-page PDF (a native LinkedIn document).

    LinkedIn's document/PDF format is a different feed object than a multi-image post —
    it renders as a swipeable, downloadable deck. The slides are the same 1080x1080 PNGs
    ``create_carousel_slide_images`` produces, so a document post is a carousel published
    through the document path.

    Returns the absolute PDF path, or None when no slide image could be read (never a
    partial/placeholder deck — the caller flags the post 'error' instead).
    """
    from PIL import Image

    from cqc_lem.utilities.logger import log_warning

    pages = []
    for path in image_paths or []:
        try:
            with Image.open(path) as img:
                pages.append(img.convert("RGB"))
        except Exception as e:
            log_warning("Could not read carousel slide for PDF", exc=e, post_id=post_id)
            return None

    if not pages:
        return None

    if output_dir is None:
        current_dir = os.path.dirname(__file__)
        assets_root = os.path.join(current_dir, "..", "assets", "images", "carousel", str(post_id))
        output_dir = os.path.realpath(assets_root)
    os.makedirs(output_dir, exist_ok=True)

    pdf_path = os.path.join(output_dir, f"document_{post_id}.pdf")
    pages[0].save(pdf_path, "PDF", save_all=True, append_images=pages[1:], resolution=150.0)
    return pdf_path


if __name__ == "__main__":
    # Debug
    # debug_master_slide_placeholders_and_text()

    # test_caption_only_slide()
    # test_create_educational_ppt()
    test_create_case_study_ppt()

    exit(0)

    # Example usage of the TestimonialCarousel model
    carousel_data = {
        "cover": {
            "title": "What Our Clients Are Saying",
            "content": "Hear from our satisfied clients about their experience "
        },
        "testimonials": [
            {
                "title": "Client Testimonial 1",
                "content": "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua.",
                "image_url": "https://example.com/client1.jpg",
                "client_name": "John Doe",
                "client_logo_url": "https://example.com/logo1.jpg"
            },
            {
                "title": "Client Testimonial 2",
                "content": "Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris nisi ut aliquip ex ea commodo consequat.",
                "image_url": "https://example.com/client2.jpg",
                "client_name": "Jane Smith",
                "client_logo_url": "https://example.com/logo2.jpg"
            }
        ],
        "call_to_action": {
            "title": "Contact Us",
            "content": "If you are ready tto learn more about our services."
        }

    }
    ppt = create_ppt("Testimonials", TestimonialCarousel(**carousel_data))
    print(f"Presentation created: {ppt}")
